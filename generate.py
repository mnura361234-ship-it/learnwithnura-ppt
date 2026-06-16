import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches

from assets_manager import ensure_output_directory
from config import LAYOUT_ORDER, OUTPUT_DIR, OUTPUT_FILENAME_TEMPLATE, SLIDE_HEIGHT, SLIDE_WIDTH, TOPIC
from developer_tools import project_health_check
from icon_manager import add_icon, get_slide_icon, get_topic_icon
from image_manager import select_topic_image
from layouts.registry import get_layout_class
from parser import load_slides
from themes import select_theme
from utils import add_background, add_topic_image, create_textbox

OUTPUT_FILENAME = OUTPUT_FILENAME_TEMPLATE.format(topic=TOPIC.replace(' ', '_'))
OUTPUT_PATH = os.path.join(OUTPUT_DIR, OUTPUT_FILENAME)


def build_cover_slide(prs, context):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, context['theme'], context['slide_height'])

    create_textbox(
        slide,
        left=Inches(0.8),
        top=Inches(2.0),
        width=Inches(6.0),
        height=Inches(0.8),
        text=context['topic'],
        theme=context['theme'],
        font_size=30,
        bold=True,
    )

    create_textbox(
        slide,
        left=Inches(0.8),
        top=Inches(3.0),
        width=Inches(5.5),
        height=Inches(0.5),
        text='LearnWithNura • AI | Web3 | Blockchain | Technology',
        theme=context['theme'],
        font_size=14,
        color=context['theme'].secondary_color,
    )

    add_topic_image(
        slide,
        context['image_path'],
        left=Inches(8.3),
        top=Inches(1.5),
        width=Inches(4.2),
        height=Inches(4.2),
    )

    return slide


def choose_layout(index, total_slides, sequence):
    if index == total_slides:
        return 'call_to_action'
    return sequence[(index - 1) % len(sequence)]


def print_debug_info(context, slide_count):
    print('=' * 40)
    print('LearnWithNura PPT Generator')
    print('=' * 40)
    print(f"Topic           : {context['topic']}")
    print(f"Selected Image  : {context['image_name']}")
    print(f"Image Path      : {context['image_path']}")
    print(f"Image Exists    : {os.path.exists(context['image_path']) if context['image_path'] else False}")
    print(f"Slides Loaded   : {slide_count}")
    print(f"Output File     : {OUTPUT_PATH}")
    print(f"Layouts Enabled : {len(context['layout_sequence']) + 1}")
    print('=' * 40)


def parse_args(args=None):
    parser = argparse.ArgumentParser(
        description='LearnWithNura PPT Generator',
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    parser.add_argument(
        '--health',
        action='store_true',
        help='Run a lightweight project health check and exit.',
    )
    return parser.parse_args(args)


def main(args=None):
    if args is None:
        args = parse_args()
    elif isinstance(args, (list, tuple)):
        args = parse_args(args)

    if args.health:
        success = project_health_check()
        sys.exit(0 if success else 1)

    try:
        ensure_output_directory(OUTPUT_DIR)
    except Exception as error:
        print(f'ERROR: Unable to prepare output directory: {error}')
        sys.exit(1)

    try:
        slides = load_slides()
    except Exception as error:
        print('\nERROR: Failed to load slides.')
        print(error)
        sys.exit(1)

    theme = select_theme(TOPIC)
    image_path, image_name = select_topic_image(TOPIC)
    layout_sequence = LAYOUT_ORDER or ['standard', 'reversed', 'key_takeaway', 'comparison', 'standard', 'reversed']
    topic_icon = get_topic_icon(TOPIC, theme.icon_pack)

    context = {
        'topic': TOPIC,
        'theme': theme,
        'image_path': image_path,
        'image_name': image_name,
        'slide_width': SLIDE_WIDTH,
        'slide_height': SLIDE_HEIGHT,
        'layout_sequence': layout_sequence,
        'topic_icon': topic_icon,
        'icon_pack': theme.icon_pack,
    }

    print_debug_info(context, len(slides))

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Check if first slide is already a cover; if so, skip the hardcoded cover
    first_slide_is_cover = slides and slides[0].layout == "cover" if slides else False
    if not first_slide_is_cover:
        build_cover_slide(prs, context)

    total_slides = len(slides)
    for index, slide in enumerate(slides, start=1):
        layout_name = slide.layout if slide.layout else choose_layout(index, total_slides, layout_sequence)
        layout_class = get_layout_class(layout_name)
        if layout_class is None:
            print(f"WARNING: Unknown layout '{layout_name}', falling back to standard.")
            layout_name = 'standard'
            layout_class = get_layout_class(layout_name)

        layout = layout_class()
        context['slide_number'] = index
        context['slide_icon'] = get_slide_icon(layout_name, theme.icon_pack)

        layout.build(prs, slide.title, slide.body, context)

    try:
        prs.save(OUTPUT_PATH)
    except PermissionError:
        print('ERROR: Permission denied while saving the presentation.')
        sys.exit(1)
    except Exception as error:
        print('ERROR: Failed to save presentation.')
        print(error)
        sys.exit(1)

    print('=' * 40)
    print('Presentation created successfully!')
    print(os.path.abspath(OUTPUT_PATH))
    print('=' * 40)


if __name__ == '__main__':
    main()
