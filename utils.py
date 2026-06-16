from assets_manager import find_asset
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def safe_text(text):
    return str(text or "").strip()


def create_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    theme,
    font_size,
    bold=False,
    color=None,
    alignment=PP_ALIGN.LEFT,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.clear()

    paragraph = frame.paragraphs[0]
    paragraph.text = safe_text(text)
    paragraph.alignment = alignment

    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.name = theme.font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or theme.primary_color

    return textbox


def add_topic_image(slide, image_path, left, top, width, height):
    if not image_path:
        return
    try:
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    except Exception:
        pass


def add_footer(slide, slide_number, theme):
    create_textbox(
        slide,
        left=Inches(0.45),
        top=Inches(6.9),
        width=Inches(4.2),
        height=Inches(0.25),
        text=f"LearnWithNura  |  Slide {slide_number}",
        theme=theme,
        font_size=9,
        bold=False,
        color=theme.secondary_color,
    )


def add_background(slide, theme, slide_height):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme.background_color

    if theme.background_image:
        background_image_path = find_asset(theme.background_image)
        if background_image_path:
            try:
                slide.shapes.add_picture(
                    background_image_path,
                    0,
                    0,
                    width=slide.part.slide_layout.slide_width,
                    height=slide_height,
                )
            except Exception:
                pass

    if theme.decorative_shapes:
        for shape_info in theme.decorative_shapes:
            try:
                shape = slide.shapes.add_shape(
                    shape_info.get("shape_type", 1),
                    Inches(shape_info.get("left", 0)),
                    Inches(shape_info.get("top", 0)),
                    Inches(shape_info.get("width", 1)),
                    Inches(shape_info.get("height", 1)),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = shape_info.get("color", theme.accent_color)
                shape.line.fill.background()
            except Exception:
                continue

    bar = slide.shapes.add_shape(
        1,
        0,
        0,
        Inches(0.18),
        slide_height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.accent_color
    bar.line.fill.background()


def split_comparison(body):
    """Split comparison content using | if available."""
    if "|" in body:
        left_text, right_text = body.split("|", 1)
        return left_text.strip(), right_text.strip()

    midpoint = len(body) // 2
    return body[:midpoint].strip(), body[midpoint:].strip()


def parse_timeline(body):
    """Parse timeline entries formatted as label::description per line."""
    entries = []
    if not body:
        return entries

    for raw_line in body.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        parts = [part.strip() for part in line.split("::", 1)]
        if len(parts) != 2 or not parts[0] or not parts[1]:
            continue

        entries.append((parts[0], parts[1]))

    return entries


def draw_timeline(slide, entries, theme, slide_width=None):
    """Render a horizontal timeline with nodes and labels."""
    from components.timeline_component import TimelineComponent

    TimelineComponent().render(
        slide,
        entries,
        theme,
        slide_width,
    )
