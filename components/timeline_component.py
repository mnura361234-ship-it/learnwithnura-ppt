from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from utils import create_textbox


class TimelineComponent:
    """Render timeline visual elements on a slide.

    This component encapsulates timeline-specific drawing logic so layouts
    can remain composition-focused.
    """

    def render(self, slide, entries, theme, slide_width):
        if not entries or slide_width is None:
            return

        max_nodes = min(len(entries), 5)
        left_margin = Inches(0.8)
        right_margin = Inches(0.8)
        total_width = slide_width - left_margin - right_margin
        node_spacing = total_width / max_nodes
        line_top = Inches(3.2)

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left_margin,
            line_top,
            total_width,
            Inches(0.14),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme.secondary_color
        bar.line.fill.background()

        for index, (label, description) in enumerate(entries[:max_nodes]):
            center_x = left_margin + node_spacing * index + node_spacing / 2
            dot_left = center_x - Inches(0.16)
            dot_top = line_top - Inches(0.18)

            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                dot_left,
                dot_top,
                Inches(0.32),
                Inches(0.32),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = theme.accent_color
            dot.line.fill.background()

            create_textbox(
                slide,
                center_x - Inches(1.2),
                line_top + Inches(0.2),
                Inches(2.4),
                Inches(0.25),
                label,
                theme,
                font_size=12,
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )

            create_textbox(
                slide,
                center_x - Inches(1.4),
                line_top + Inches(0.45),
                Inches(2.8),
                Inches(0.8),
                description,
                theme,
                font_size=11,
                alignment=PP_ALIGN.CENTER,
            )
